"""
EIZ-Legaltech Personen-Dossier PowerPoint Template
====================================================
Reproduziert das Design der EIZ-Legaltech Vorlage.
Erstellt pro Person eine Folie mit Foto, Basisprofil und Recherche-Ergebnissen.

Design-Spezifikation (aus der Vorlage extrahiert):
- Foliengrösse: 13.33 x 7.50 Zoll (Widescreen 16:9)
- Primärfarbe (accent1): #156082 (dunkles Petrol/Blau)
- Sekundärfarbe (dk2): #0E2841 (sehr dunkles Blau)
- Hintergrund-Gradient: Dunkler Bereich links mit Gradient-Effekt
- Schrift Überschriften: Aptos Display (Fallback: Calibri)
- Schrift Body: Aptos (Fallback: Calibri)
- Logo: EuropaInstitut, Position oben rechts (11.85", 0.34")
- Footer: Datum links, Name|Workshop Mitte, Seitenzahl rechts

Verwendung:
    from eiz_dossier_template import create_dossier_presentation

    personen = [
        {
            "name": "Dr. Anna Muster",
            "titel": "Partnerin",
            "organisation": "Walder Wyss",
            "typ": "jurist",  # oder "unternehmen"
            "foto_pfad": "foto_anna.jpg",  # optional, None wenn kein Foto
            "standort": "Zürich",
            "werdegang": "Universität Zürich (Dr. iur.), LL.M. Harvard",
            "fachgebiete": "M&A, Gesellschaftsrecht, Kapitalmarktrecht",
            "linkedin_zusammenfassung": "Aktiv zu Themen ESG und Corporate Governance...",
            "publikationen": "Diverse Beiträge in AJP, SZW...",
            "rankings": "Legal 500: Band 2 Corporate/M&A",
            "engagement": "Vorstandsmitglied SAV Zürich",
            "gespraechspunkte": [
                "Kürzlicher Artikel zu ESG-Pflichten",
                "Vortrag an der Uni Zürich zu Aktienrechtsrevision",
                "LinkedIn-Post über KI im Recht"
            ],
            "einschaetzung": "Strukturiert und fachlich orientiert. Schätzt gut vorbereitete Gesprächspartner.",
            "zusatz_info": ""  # Freitextfeld für zusätzliche Infos
        }
    ]

    create_dossier_presentation(
        personen=personen,
        output_pfad="dossier_output.pptx",
        logo_pfad="europa_institut_logo.jpg",
        autor_name="Ramona Schindler",
        workshop_name="Legaltech Workshop"
    )
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
from datetime import datetime


# ============================================================
# DESIGN-KONSTANTEN (aus der Vorlage extrahiert)
# ============================================================

# Foliengrösse
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.50)

# Farben
COLOR_ACCENT1 = RGBColor(0x15, 0x60, 0x82)      # #156082 - Hauptfarbe (Petrol)
COLOR_DARK = RGBColor(0x0E, 0x28, 0x41)           # #0E2841 - Sehr dunkles Blau
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)           # Schwarz
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)           # Weiss
COLOR_LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)     # #E8E8E8 - Helles Grau
COLOR_MID_GRAY = RGBColor(0x88, 0x88, 0x88)        # Mittleres Grau
COLOR_FOOTER_TEXT = RGBColor(0x88, 0x88, 0x88)      # Footer-Textfarbe

# Schriften (Aptos mit Fallback auf Calibri)
FONT_HEADING = "Aptos Display"
FONT_BODY = "Aptos"
FONT_FALLBACK = "Calibri"

# Positionen
LOGO_LEFT = Inches(11.85)
LOGO_TOP = Inches(0.34)
LOGO_WIDTH = Inches(1.12)
LOGO_HEIGHT = Inches(0.61)

# Linkes Panel (dunkler Bereich)
LEFT_PANEL_WIDTH = Inches(4.0)

# Footer
FOOTER_TOP = Inches(6.95)
FOOTER_HEIGHT = Inches(0.40)


# ============================================================
# HILFSFUNKTIONEN
# ============================================================

def _add_dark_left_panel(slide):
    """Erstellt den dunklen linken Hintergrundbereich (wie in der Vorlage)."""
    # Vollflächiger Hintergrund (leichtes Grau/Weiss)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_WHITE
    bg.line.fill.background()

    # Dunkles linkes Panel
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(-0.5), Inches(0), Inches(5.0), SLIDE_HEIGHT
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLOR_DARK
    panel.line.fill.background()

    # Overlay für Gradient-Effekt (halbtransparenter Kreis wie in der Vorlage)
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-1.0), Inches(1.0), Inches(5.5), Inches(5.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLOR_ACCENT1
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    srgb = overlay._element.find('.//' + qn('a:srgbClr'))
    if srgb is not None:
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', '30000')
    overlay.line.fill.background()


def _add_logo(slide, logo_pfad):
    """Fügt das Logo oben rechts ein."""
    if logo_pfad and os.path.exists(logo_pfad):
        slide.shapes.add_picture(
            logo_pfad, LOGO_LEFT, LOGO_TOP, LOGO_WIDTH, LOGO_HEIGHT
        )


def _add_footer(slide, datum, autor_name, workshop_name, slide_nummer):
    """Fügt den Footer ein: Datum links, Name|Workshop Mitte, Seitenzahl rechts."""
    # Trennlinie
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    line.line.fill.background()

    # Datum links
    txBox = slide.shapes.add_textbox(Inches(0.92), FOOTER_TOP, Inches(3.0), FOOTER_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = datum
    run.font.size = Pt(10)
    run.font.name = FONT_BODY
    run.font.color.rgb = COLOR_FOOTER_TEXT

    # Name | Workshop Mitte
    txBox = slide.shapes.add_textbox(Inches(4.42), FOOTER_TOP, Inches(4.50), FOOTER_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{autor_name} | {workshop_name}"
    run.font.size = Pt(10)
    run.font.name = FONT_BODY
    run.font.color.rgb = COLOR_FOOTER_TEXT

    # Seitenzahl rechts
    txBox = slide.shapes.add_textbox(Inches(9.42), FOOTER_TOP, Inches(3.0), FOOTER_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(slide_nummer)
    run.font.size = Pt(10)
    run.font.name = FONT_BODY
    run.font.color.rgb = COLOR_FOOTER_TEXT


def _add_text_run(paragraph, text, font_name=None, font_size=None,
                  bold=None, italic=None, color=None):
    """Fügt einen formatierten Text-Run zu einem Paragraphen hinzu."""
    run = paragraph.add_run()
    run.text = text
    if font_name:
        run.font.name = font_name
    if font_size:
        run.font.size = font_size
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return run


# ============================================================
# TITELFOLIE
# ============================================================

def _create_title_slide(prs, titel, untertitel, autor_name, datum, logo_pfad):
    """Erstellt die Titelfolie im EIZ-Design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Weisser Hintergrund oben
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_WHITE
    bg.line.fill.background()

    # Dunkler Balken unten (wie Titelfolie der Vorlage)
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(4.78), SLIDE_WIDTH, Inches(2.72)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = COLOR_DARK
    bottom_bar.line.fill.background()

    # Petrol-Overlay auf dem dunklen Balken
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.90), Inches(4.78), Inches(4.43), Inches(2.72)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLOR_ACCENT1
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    srgb = overlay._element.find('.//' + qn('a:srgbClr'))
    if srgb is not None:
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', '50000')
    overlay.line.fill.background()

    # Titel
    txBox = slide.shapes.add_textbox(
        Inches(1.45), Inches(1.25), Inches(9.86), Inches(2.63)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _add_text_run(p, titel, FONT_HEADING, Pt(48), color=COLOR_BLACK)

    # Untertitel (falls vorhanden)
    if untertitel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(6)
        _add_text_run(p2, untertitel, FONT_HEADING, Pt(28), color=COLOR_MID_GRAY)

    # Autor und Datum im dunklen Balken
    txBox = slide.shapes.add_textbox(
        Inches(1.45), Inches(5.35), Inches(6.99), Inches(1.18)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _add_text_run(p, autor_name, FONT_BODY, Pt(22), color=COLOR_WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    _add_text_run(p2, datum, FONT_BODY, Pt(22), color=COLOR_WHITE)

    # Logo
    _add_logo(slide, logo_pfad)


# ============================================================
# PERSONEN-DOSSIER-FOLIE
# ============================================================

def _create_person_slide(prs, person, slide_nummer, logo_pfad,
                         autor_name, workshop_name, datum):
    """
    Erstellt eine Dossier-Folie für eine Person.

    Layout:
    - Links (dunkel): Foto + Name + Titel + Organisation
    - Rechts (hell): Recherche-Ergebnisse in kompaktem Layout
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Hintergrund
    _add_dark_left_panel(slide)

    # --- LINKES PANEL: Foto + Name ---

    # Foto einfügen
    foto_pfad = person.get("foto_pfad")
    if foto_pfad and os.path.exists(foto_pfad):
        # Foto als Kreis-ähnliches Element (rechteckig mit Position)
        pic = slide.shapes.add_picture(
            foto_pfad,
            Inches(0.60), Inches(0.90),
            Inches(2.60), Inches(2.60)
        )
    else:
        # Platzhalter-Kreis für fehlendes Foto
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.90), Inches(0.90),
            Inches(2.00), Inches(2.00)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = COLOR_ACCENT1
        placeholder.line.fill.background()

        # "Kein Foto" Text
        txBox = slide.shapes.add_textbox(
            Inches(0.90), Inches(1.60), Inches(2.00), Inches(0.60)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _add_text_run(p, "Kein Foto", FONT_BODY, Pt(12), color=COLOR_WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        _add_text_run(p2, "verfügbar", FONT_BODY, Pt(12), color=COLOR_WHITE)

    # Name
    txBox = slide.shapes.add_textbox(
        Inches(0.30), Inches(3.70), Inches(3.50), Inches(1.20)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _add_text_run(p, person.get("name", ""), FONT_HEADING, Pt(28), bold=True, color=COLOR_WHITE)

    # Titel + Organisation
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.RIGHT
    p2.space_before = Pt(4)
    titel_org = person.get("titel", "")
    if person.get("organisation"):
        titel_org += f"\n{person['organisation']}"
    _add_text_run(p2, titel_org, FONT_BODY, Pt(16), color=COLOR_WHITE)

    # Standort
    if person.get("standort"):
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.RIGHT
        p3.space_before = Pt(4)
        _add_text_run(p3, person["standort"], FONT_BODY, Pt(14), italic=True, color=COLOR_LIGHT_GRAY)

    # Typ-Indikator
    typ_text = "Jurist/in" if person.get("typ") == "jurist" else "Unternehmen"
    txBox = slide.shapes.add_textbox(
        Inches(0.30), Inches(5.80), Inches(3.50), Inches(0.35)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    _add_text_run(p, typ_text, FONT_BODY, Pt(11), italic=True, color=COLOR_ACCENT1)

    # --- RECHTES PANEL: Recherche-Ergebnisse ---

    right_left = Inches(4.50)
    right_width = Inches(8.30)
    current_top = Inches(0.40)
    line_height = Pt(13)

    def add_section(label, content, top_pos):
        """Fügt einen Abschnitt (Label + Inhalt) hinzu und gibt die neue Y-Position zurück."""
        if not content:
            return top_pos

        # Label
        txBox = slide.shapes.add_textbox(right_left, top_pos, right_width, Inches(0.30))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_text_run(p, label, FONT_HEADING, Pt(13), bold=True, color=COLOR_ACCENT1)
        top_pos += Inches(0.28)

        # Inhalt
        txBox = slide.shapes.add_textbox(right_left, top_pos, right_width, Inches(0.80))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_text_run(p, str(content), FONT_BODY, Pt(11), color=COLOR_BLACK)
        # Höhe basierend auf Textlänge schätzen
        lines = max(1, len(str(content)) // 90 + str(content).count('\n') + 1)
        top_pos += Inches(0.18 * lines + 0.08)

        return top_pos

    # Werdegang
    current_top = add_section("Werdegang", person.get("werdegang"), current_top)

    # Fachgebiete / Verantwortungsbereich
    if person.get("typ") == "jurist":
        current_top = add_section("Rechtsgebiete", person.get("fachgebiete"), current_top)
    else:
        current_top = add_section("Verantwortungsbereich", person.get("fachgebiete"), current_top)

    # LinkedIn
    current_top = add_section("LinkedIn & Online-Präsenz", person.get("linkedin_zusammenfassung"), current_top)

    # Publikationen
    current_top = add_section("Publikationen & Medien", person.get("publikationen"), current_top)

    # Rankings (nur bei Juristen)
    if person.get("typ") == "jurist" and person.get("rankings"):
        current_top = add_section("Rankings", person.get("rankings"), current_top)

    # Engagement
    current_top = add_section("Engagement & Netzwerk", person.get("engagement"), current_top)

    # Zusatzinfo (z.B. Unternehmensprofil)
    if person.get("zusatz_info"):
        current_top = add_section("Weitere Informationen", person.get("zusatz_info"), current_top)

    # --- GESPRÄCHSVORBEREITUNG (unterer Bereich rechts) ---

    # Trennlinie
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        right_left, current_top, right_width, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT1
    line.line.fill.background()
    current_top += Inches(0.15)

    # Gesprächsanknüpfungspunkte
    if person.get("gespraechspunkte"):
        txBox = slide.shapes.add_textbox(right_left, current_top, right_width, Inches(0.25))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_text_run(p, "Gesprächsanknüpfungspunkte", FONT_HEADING, Pt(13), bold=True, color=COLOR_ACCENT1)
        current_top += Inches(0.25)

        for punkt in person["gespraechspunkte"]:
            txBox = slide.shapes.add_textbox(
                right_left + Inches(0.15), current_top, right_width - Inches(0.15), Inches(0.22)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            _add_text_run(p, f"→ {punkt}", FONT_BODY, Pt(11), color=COLOR_BLACK)
            current_top += Inches(0.22)

    # Einschätzung
    if person.get("einschaetzung"):
        current_top += Inches(0.05)
        txBox = slide.shapes.add_textbox(right_left, current_top, right_width, Inches(0.25))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_text_run(p, "Einschätzung", FONT_HEADING, Pt(13), bold=True, color=COLOR_ACCENT1)
        current_top += Inches(0.25)

        txBox = slide.shapes.add_textbox(right_left, current_top, right_width, Inches(0.50))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_text_run(p, person["einschaetzung"], FONT_BODY, Pt(11), italic=True, color=COLOR_MID_GRAY)

    # Logo
    _add_logo(slide, logo_pfad)

    # Footer
    _add_footer(slide, datum, autor_name, workshop_name, slide_nummer)


# ============================================================
# HAUPTFUNKTION
# ============================================================

def create_dossier_presentation(
    personen,
    output_pfad="dossier_output.pptx",
    logo_pfad=None,
    autor_name="",
    workshop_name="",
    titel="Personen-Dossier",
    untertitel="Gesprächsvorbereitung",
    datum=None
):
    """
    Erstellt eine komplette Dossier-Präsentation.

    Args:
        personen: Liste von Personen-Dicts (siehe Modul-Docstring für Struktur)
        output_pfad: Pfad für die Ausgabe-PPTX
        logo_pfad: Pfad zum Logo (JPG/PNG)
        autor_name: Name des Autors (für Footer + Titelfolie)
        workshop_name: Name des Workshops (für Footer)
        titel: Titel der Präsentation
        untertitel: Untertitel der Präsentation
        datum: Datum-String (Standard: heutiges Datum)

    Returns:
        Pfad zur erstellten PPTX-Datei
    """
    if datum is None:
        datum = datetime.now().strftime("%d.%m.%Y")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Titelfolie
    _create_title_slide(prs, titel, untertitel, autor_name, datum, logo_pfad)

    # Eine Folie pro Person
    for i, person in enumerate(personen):
        _create_person_slide(
            prs, person,
            slide_nummer=i + 2,  # Titelfolie ist 1
            logo_pfad=logo_pfad,
            autor_name=autor_name,
            workshop_name=workshop_name,
            datum=datum
        )

    prs.save(output_pfad)
    return output_pfad


# ============================================================
# DEMO / TEST
# ============================================================

if __name__ == "__main__":
    # Demo mit Beispieldaten
    demo_personen = [
        {
            "name": "Dr. Anna Muster",
            "titel": "Partnerin",
            "organisation": "Walder Wyss",
            "typ": "jurist",
            "foto_pfad": None,
            "standort": "Zürich",
            "werdegang": "Dr. iur. Universität Zürich (2012), LL.M. Harvard Law School (2014). Zuvor Associate bei Lenz & Staehelin (2008–2015).",
            "fachgebiete": "M&A, Gesellschaftsrecht, Kapitalmarktrecht, Private Equity",
            "linkedin_zusammenfassung": "Aktiv auf LinkedIn mit Posts zu ESG-Regulierung und Corporate Governance. Regelmässige Kommentare zu Aktienrechtsrevision.",
            "publikationen": "Co-Autorin Kommentar zum Aktienrecht (Basler Kommentar, 2023). Diverse Beiträge in AJP und SZW. Vortrag an der IBA Konferenz 2024.",
            "rankings": "Legal 500 2024: Band 2 Corporate/M&A. Chambers Europe: Up and Coming.",
            "engagement": "Vorstandsmitglied SAV Sektion Zürich. Lehrauftrag Universität Zürich (Gesellschaftsrecht).",
            "gespraechspunkte": [
                "Kürzlicher AJP-Artikel zu ESG-Pflichten des Verwaltungsrats",
                "Lehrauftrag an der Uni Zürich – Erfahrungen mit Studierenden",
                "LinkedIn-Diskussion über KI-Tools in der Vertragsanalyse"
            ],
            "einschaetzung": "Strukturiert und fachlich orientiert. Schätzt gut vorbereitete Gesprächspartner. Offen für Innovation im Legal-Bereich.",
            "zusatz_info": ""
        },
        {
            "name": "Thomas Beispiel",
            "titel": "CFO",
            "organisation": "Helvetia Versicherungen",
            "typ": "unternehmen",
            "foto_pfad": None,
            "standort": "St. Gallen / Basel",
            "werdegang": "Betriebswirtschaft HSG (2005), CFA Charterholder. Zuvor Head of Finance bei Zurich Insurance (2010–2018).",
            "fachgebiete": "Finanzsteuerung, Risikomanagement, Investor Relations, M&A-Integration",
            "linkedin_zusammenfassung": "Regelmässige Posts zu Versicherungstrends, Nachhaltigkeit in der Finanzbranche und Leadership.",
            "publikationen": "Interview in der NZZ zum Thema Solvency II (März 2025). Panelist am Swiss Insurance Forum.",
            "rankings": "",
            "engagement": "Mitglied Finanzkommission Economiesuisse. Verwaltungsrat eines Fintech-Startups.",
            "gespraechspunkte": [
                "Helvetias Nachhaltigkeitsstrategie und Impact Investing",
                "NZZ-Interview zu regulatorischen Herausforderungen",
                "Verwaltungsratsmandat beim Fintech – Interesse an Innovation"
            ],
            "einschaetzung": "Zahlenorientiert und pragmatisch. Erwartet klare, faktenbasierte Argumente. Gleichzeitig innovationsaffin.",
            "zusatz_info": "Helvetia Versicherungen: Börsenkotiert (SIX), ca. 12'000 MA, Prämienvolumen CHF 11 Mrd. Aktuell Fokus auf Digitalisierung und Expansion in Specialty Lines."
        }
    ]

    output = create_dossier_presentation(
        personen=demo_personen,
        output_pfad="/home/claude/demo_dossier.pptx",
        logo_pfad="/home/claude/images/Grafik_7.jpg",  # Logo aus der Vorlage
        autor_name="Ramona Schindler",
        workshop_name="Legaltech Workshop",
        titel="Personen-Dossier",
        untertitel="Gesprächsvorbereitung – Meeting 15.03.2026"
    )
    print(f"Erstellt: {output}")
