"""
Generiert 50 realistische, chaotisch benannte Dateien für den Workshop-Demo.
Mix aus: Verträgen, Urteilen, Fachartikeln, Korrespondenz, Rechnungen,
Präsentationen, Notizen – in DE, FR, EN.
"""

import os
from docx import Document
from docx.shared import Pt, Inches
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from datetime import datetime, timedelta
import random

OUT = "/home/claude/workshop-dateien"
os.makedirs(OUT, exist_ok=True)


def make_docx(filename, title, paragraphs):
    doc = Document()
    doc.add_heading(title, level=1)
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save(os.path.join(OUT, filename))


def make_pdf(filename, title, paragraphs):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    # Use built-in font with latin-1 encoding
    pdf.set_font("Helvetica", "B", 16)
    safe_title = title.encode('latin-1', 'replace').decode('latin-1')
    pdf.cell(0, 10, safe_title, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 11)
    for p in paragraphs:
        safe_p = p.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 6, safe_p)
        pdf.ln(3)
    pdf.output(os.path.join(OUT, filename))


def make_txt(filename, content):
    with open(os.path.join(OUT, filename), "w", encoding="utf-8") as f:
        f.write(content)


def make_pptx(filename, title, slides_content):
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = "Präsentation"
    for stitle, sbody in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = stitle
        slide.placeholders[1].text = sbody
    prs.save(os.path.join(OUT, filename))


# ============================================================
# VERTRÄGE (10 Dateien)
# ============================================================

make_docx("NDA_Project_Alpha_unsigned.docx",
    "Geheimhaltungsvereinbarung (NDA)",
    [
        "zwischen der TechVision AG, Zürich (nachfolgend «Offenlegende Partei») und der InnoSoft GmbH, München (nachfolgend «Empfangende Partei»)",
        "1. Gegenstand: Die Parteien beabsichtigen, vertrauliche Informationen im Zusammenhang mit dem Projekt Alpha auszutauschen.",
        "2. Vertrauliche Informationen: Als vertraulich gelten sämtliche Informationen technischer, wirtschaftlicher oder organisatorischer Natur, die im Rahmen der Zusammenarbeit offengelegt werden.",
        "3. Geheimhaltungspflicht: Die empfangende Partei verpflichtet sich, die vertraulichen Informationen streng geheim zu halten und nur für den vereinbarten Zweck zu verwenden.",
        "4. Dauer: Diese Vereinbarung tritt mit Unterzeichnung in Kraft und gilt für einen Zeitraum von 3 Jahren.",
        "5. Vertragsstrafe: Bei Verletzung der Geheimhaltungspflicht wird eine Konventionalstrafe von CHF 100'000 fällig.",
        "6. Gerichtsstand: Zürich. Anwendbares Recht: Schweizerisches Recht.",
        "\n\nOrt, Datum: _________________\n\nTechVision AG: _________________\n\nInnoSoft GmbH: _________________"
    ])

make_docx("Vertrag_Meier_v3_FINAL.docx",
    "Dienstleistungsvertrag",
    [
        "zwischen Meier Consulting AG, Bern und DataFlow Systems AG, Zürich",
        "Vertrags-Nr.: DV-2025-0847",
        "1. Vertragsgegenstand: Meier Consulting erbringt Beratungsdienstleistungen im Bereich Datenschutz-Compliance gemäss dem revidierten Datenschutzgesetz (nDSG).",
        "2. Leistungsumfang: a) Gap-Analyse der bestehenden Datenschutzmassnahmen, b) Erarbeitung eines Massnahmenplans, c) Schulung der Mitarbeitenden, d) Begleitung der Implementierung.",
        "3. Vergütung: Stundenhonorar CHF 350.– zzgl. MWST. Maximaler Aufwand: 200 Stunden. Gesamtbudget: CHF 70'000.– zzgl. MWST.",
        "4. Laufzeit: 01.01.2026 bis 31.12.2026. Kündigungsfrist: 3 Monate zum Quartalsende.",
        "5. Haftung: Die Haftung von Meier Consulting ist auf den Vertragswert beschränkt.",
        "6. Gerichtsstand: Bern."
    ])

make_pdf("arbeitsvertrag_vorlage_2024.pdf",
    "Arbeitsvertrag (Muster)",
    [
        "zwischen [Arbeitgeber] und [Arbeitnehmer/in]",
        "1. Beginn und Dauer: Das Arbeitsverhaeltnis beginnt am [Datum] und wird auf unbestimmte Zeit abgeschlossen.",
        "2. Funktion: [Funktionsbezeichnung], Pensum: [100%]",
        "3. Arbeitsort: [Ort]",
        "4. Probezeit: 3 Monate",
        "5. Kuendigungsfrist: Im 1. Dienstjahr 1 Monat, ab 2. Dienstjahr 2 Monate, ab 10. Dienstjahr 3 Monate.",
        "6. Verguetung: CHF [Betrag] brutto pro Jahr, zahlbar in 13 Monatsloeehnen.",
        "7. Ferien: 25 Arbeitstage pro Kalenderjahr.",
        "8. Gerichtsstand: [Ort]. Anwendbares Recht: Schweizerisches Recht (OR Art. 319 ff.)"
    ])

make_docx("Mietvertrag_Büro_Talstrasse.docx",
    "Mietvertrag für Geschäftsräume",
    [
        "Vermieterin: Immobilien Zürich AG, Bahnhofstrasse 42, 8001 Zürich",
        "Mieterin: Kanzlei Weber & Partner, Talstrasse 18, 8001 Zürich",
        "Mietobjekt: Büroräumlichkeiten im 3. OG, ca. 240 m², 8 Büros, 1 Sitzungszimmer",
        "Mietzins: CHF 8'400.– pro Monat zzgl. Nebenkosten CHF 1'200.–",
        "Mietbeginn: 01.04.2025. Mindestmietdauer: 5 Jahre.",
        "Kündigungsfrist: 6 Monate auf Ende eines Kalenderquartals.",
        "Kaution: CHF 25'200.– (3 Monatsmieten), hinterlegt auf Sperrkonto.",
        "Gerichtsstand: Zürich."
    ])

make_pdf("SLA_CloudProvider_2025.pdf",
    "Service Level Agreement",
    [
        "between CloudSecure AG, Zurich and DataFlow Systems AG, Zurich",
        "Service: Managed Cloud Hosting for Legal Document Management System",
        "1. Service Availability: 99.95% uptime per calendar month.",
        "2. Response Times: Critical incidents: 15 minutes. High: 1 hour. Medium: 4 hours. Low: next business day.",
        "3. Data Location: All data stored exclusively in Swiss data centers (Zurich, Geneva).",
        "4. Backup: Daily incremental, weekly full backup. Retention: 90 days.",
        "5. Security: ISO 27001 certified. Annual penetration testing.",
        "6. Penalties: Below 99.9%: 10% credit. Below 99.5%: 25% credit. Below 99.0%: 50% credit.",
        "7. Term: 36 months. Automatic renewal for 12 months unless terminated 6 months prior.",
        "8. Governing Law: Swiss Law. Jurisdiction: Zurich."
    ])

make_docx("Lizenzvertrag_Software_Draft2.docx",
    "Softwarelizenzvertrag – ENTWURF v2",
    [
        "Lizenzgeberin: LegalTech Solutions SA, Genf",
        "Lizenznehmerin: Anwaltskanzlei Brunner & Meier, Basel",
        "Software: «JurisAssist Pro» – KI-gestützte Rechtsrecherche-Plattform",
        "1. Lizenzumfang: Nicht-exklusive Nutzungslizenz für max. 25 gleichzeitige Nutzer.",
        "2. Lizenzgebühr: CHF 2'500.– pro Monat. Jährliche Preisanpassung max. 5%.",
        "3. Laufzeit: 24 Monate ab Aktivierung. Kündigungsfrist: 6 Monate.",
        "4. Datenschutz: Verarbeitung ausschliesslich in der Schweiz. Auftragsverarbeitungsvertrag als Annex A.",
        "5. Gewährleistung: Lizenzgeberin gewährleistet Funktionalität gemäss Spezifikation.",
        "6. Haftung: Beschränkt auf den jährlichen Lizenzwert. Ausschluss von Folgeschäden.",
        "ACHTUNG: Noch nicht freigegeben – Kommentare von Partner Brunner ausstehend!"
    ])

make_pdf("contrat_bail_geneve.pdf",
    "Contrat de bail a loyer",
    [
        "Bailleur: Regie Immobiliere SA, Geneve",
        "Locataire: Etude d'avocats Martin & Associes, Geneve",
        "Objet: Locaux commerciaux, Rue du Rhone 45, 1204 Geneve, 3eme etage",
        "Loyer mensuel: CHF 12'500.-- charges comprises",
        "Duree: Du 1er janvier 2025 au 31 decembre 2029",
        "Resiliation: Preavis de 6 mois pour la fin d'un trimestre",
        "Depot de garantie: CHF 37'500.-- (3 mois de loyer)",
        "Droit applicable: Droit suisse. For juridique: Geneve."
    ])

make_docx("AGB_TechVision_v4.docx",
    "Allgemeine Geschäftsbedingungen (AGB)",
    [
        "TechVision AG, Zürich – Stand: Januar 2026",
        "1. Geltungsbereich: Diese AGB gelten für sämtliche Geschäftsbeziehungen zwischen der TechVision AG und ihren Kunden.",
        "2. Vertragsschluss: Angebote sind freibleibend. Der Vertrag kommt durch schriftliche Auftragsbestätigung zustande.",
        "3. Preise: Alle Preise verstehen sich in CHF, exkl. MWST.",
        "4. Zahlung: Zahlbar innert 30 Tagen netto.",
        "5. Gewährleistung: 24 Monate ab Lieferung.",
        "6. Haftung: Die Haftung ist auf Vorsatz und grobe Fahrlässigkeit beschränkt. Die Haftung für Folgeschäden ist ausgeschlossen.",
        "7. Datenschutz: Es gilt die Datenschutzerklärung der TechVision AG.",
        "8. Gerichtsstand: Zürich. Anwendbares Recht: Schweizerisches Recht.",
        "9. Salvatorische Klausel: Sollte eine Bestimmung unwirksam sein, bleibt die Wirksamkeit der übrigen Bestimmungen unberührt."
    ])

make_pdf("Kaufvertrag_Geschaeftsanteile_VERTRAULICH.pdf",
    "Share Purchase Agreement (Kaufvertrag Geschaeftsanteile)",
    [
        "VERTRAULICH / CONFIDENTIAL",
        "Seller: Founding Partners Holding AG, Zug",
        "Buyer: European Growth Capital Fund II LP, Luxembourg",
        "Target: MedTech Innovations AG, Zurich",
        "Purchase Price: CHF 45'000'000 (enterprise value) subject to customary adjustments",
        "Closing Conditions: (i) FINMA approval, (ii) Competition clearance, (iii) Key employee retention",
        "Representations & Warranties: Standard R&W package, 24-month survival period",
        "Indemnification Cap: 30% of purchase price",
        "Non-Compete: Sellers bound for 3 years post-closing",
        "Governing Law: Swiss Law. Arbitration: ICC Zurich.",
        "DRAFT - FOR DISCUSSION PURPOSES ONLY - 15.02.2026"
    ])

make_docx("Aufhebungsvereinbarung_Müller.docx",
    "Aufhebungsvereinbarung",
    [
        "zwischen der DataFlow Systems AG (Arbeitgeberin) und Herrn Thomas Müller (Arbeitnehmer)",
        "Die Parteien vereinbaren einvernehmlich:",
        "1. Das Arbeitsverhältnis wird per 30.06.2026 aufgelöst.",
        "2. Der Arbeitnehmer wird per sofort freigestellt. Die Freistellung wird an die Kündigungsfrist angerechnet.",
        "3. Die Arbeitgeberin zahlt eine Abfindung von CHF 85'000.– brutto.",
        "4. Der Arbeitnehmer erhält ein sehr gutes Arbeitszeugnis (Note: «stets zu unserer vollsten Zufriedenheit»).",
        "5. Der Arbeitnehmer verzichtet auf die Geltendmachung weiterer Ansprüche.",
        "6. Über den Inhalt dieser Vereinbarung wird Stillschweigen vereinbart.",
        "Zürich, den _______________"
    ])


# ============================================================
# URTEILE & BGE (8 Dateien)
# ============================================================

make_pdf("BGE_148_II_137.pdf",
    "BGE 148 II 137",
    [
        "Bundesgericht, II. zivilrechtliche Abteilung",
        "Urteil vom 22. Maerz 2022",
        "Regeste:",
        "Art. 697 OR. Auskunftsrecht des Aktionaers.",
        "Voraussetzungen und Grenzen des Auskunftsrechts des Aktionaers gegenueber der Gesellschaft. Verhaeltnis zum Einsichtsrecht (E. 3).",
        "Die Gesellschaft kann die Auskunft verweigern, soweit Geschaeftsgeheimnisse oder andere schutzwuerdige Interessen der Gesellschaft gefaehrdet wuerden (E. 4).",
        "Das Auskunftsrecht dient der Ausuebung der Aktionaersrechte und darf nicht zu sachfremden Zwecken missbraucht werden (E. 5).",
        "Im vorliegenden Fall hat die Vorinstanz das Auskunftsrecht zu Recht bejaht, da die verlangten Informationen fuer die Beurteilung der Geschaeftsfuehrung erforderlich waren (E. 6)."
    ])

make_pdf("Urteil_HGer_ZH_2025_1247.pdf",
    "Handelsgericht Zuerich - Urteil HG250124",
    [
        "Handelsgericht des Kantons Zuerich, Einzelgericht",
        "Urteil vom 15. Januar 2026",
        "Klaegerin: TechVision AG, Zuerich / Beklagte: InnoSoft GmbH, Muenchen",
        "Streitgegenstand: Verletzung der Geheimhaltungsvereinbarung (NDA) und Schadenersatzforderung",
        "Streitwert: CHF 500'000.--",
        "Erwägungen:",
        "1. Die Klaegerin macht geltend, die Beklagte habe vertrauliche Informationen an Dritte weitergegeben.",
        "2. Die Beklagte bestreitet die Weitergabe und beruft sich auf die Ausnahme fuer oeffentlich zugaengliche Informationen.",
        "3. Das Gericht stuetzt sich auf die Beweiserhebung und kommt zum Schluss, dass eine Verletzung der NDA vorliegt.",
        "Urteilsspruch: Die Beklagte wird verpflichtet, der Klaegerin CHF 150'000.-- zu bezahlen."
    ])

make_pdf("arret_TF_4A_251_2024.pdf",
    "Arret du Tribunal federal 4A_251/2024",
    [
        "Tribunal federal, Ire Cour de droit civil",
        "Arret du 18 novembre 2024",
        "Recours en matiere civile contre l'arret de la Cour de justice du canton de Geneve",
        "Objet: Responsabilite contractuelle, clause penale",
        "Regeste: Art. 160 ss CO. Clause penale excessive. Pouvoir de reduction du juge.",
        "Le Tribunal federal rappelle que le juge dispose d'un large pouvoir d'appreciation pour reduire une clause penale excessive (art. 163 al. 3 CO).",
        "En l'espece, une clause penale de CHF 500'000 pour la violation d'une clause de non-concurrence d'une duree de 2 ans est jugee excessive et reduite a CHF 200'000.",
        "Le recours est partiellement admis."
    ])

make_docx("Entscheid_EDÖB_2025.docx",
    "Empfehlung des EDÖB gemäss Art. 49 DSG",
    [
        "Eidgenössischer Datenschutz- und Öffentlichkeitsbeauftragter (EDÖB)",
        "Sachverhaltsabklärung betreffend die Nutzung von KI-gestützten Tools in Schweizer Anwaltskanzleien",
        "I. Sachverhalt: Der EDÖB hat Kenntnis erhalten, dass mehrere Anwaltskanzleien KI-Tools einsetzen, bei denen Mandantendaten in Cloud-Dienste übermittelt werden.",
        "II. Rechtliche Beurteilung: Die Übermittlung von Personendaten an ausländische Cloud-Anbieter stellt eine Bekanntgabe ins Ausland dar (Art. 16 DSG).",
        "III. Empfehlungen: 1) Durchführung einer DSFA vor Einsatz von KI-Tools. 2) Sicherstellung, dass keine Mandantengeheimnisse (Art. 13 StGB) preisgegeben werden. 3) Bevorzugung von Anbietern mit Datenhaltung in der Schweiz.",
        "IV. Frist zur Stellungnahme: 60 Tage."
    ])

make_pdf("BGer_5A_892_2023.pdf",
    "Bundesgericht 5A_892/2023",
    [
        "Bundesgericht, II. zivilrechtliche Abteilung",
        "Urteil vom 4. Juni 2024",
        "Beschwerde gegen das Urteil des Obergerichts des Kantons Bern",
        "Gegenstand: Erbteilung, Herabsetzungsklage",
        "Art. 522 ff. ZGB: Pflichtteile und Herabsetzung",
        "Regeste: Der Pflichtteilsschutz bezweckt, den gesetzlichen Erben einen Mindestanteil am Nachlass zu sichern. Verfuegungen von Todes wegen, die den Pflichtteil verletzen, sind auf Klage hin herabzusetzen.",
        "Im vorliegenden Fall hat der Erblasser den gesamten Nachlass einem gemeinnuetzigen Verein vermacht und seine drei Kinder enterbt.",
        "Die Beschwerde wird gutgeheissen. Die lebzeitige Schenkung wird im Umfang der Pflichtteilsverletzung herabgesetzt."
    ])

make_pdf("CJUE_C-311_18_Schrems_II.pdf",
    "CJEU Case C-311/18 (Schrems II) - Summary",
    [
        "Court of Justice of the European Union",
        "Judgment of 16 July 2020 - Case C-311/18",
        "Data Protection Commissioner v Facebook Ireland and Maximillian Schrems",
        "Key Holdings:",
        "1. The EU-US Privacy Shield Decision is invalid.",
        "2. Standard Contractual Clauses (SCCs) remain valid in principle.",
        "3. Data exporters must verify that the level of protection in the third country is essentially equivalent to that in the EU.",
        "4. Supplementary measures may be required where SCCs alone are insufficient.",
        "Impact: This judgment fundamentally changed international data transfer practices and led to the development of the EU-US Data Privacy Framework.",
        "Note: This summary is for reference purposes only."
    ])

make_pdf("BVGE_2024_IV_3.pdf",
    "Bundesverwaltungsgericht BVGE 2024/IV/3",
    [
        "Bundesverwaltungsgericht, Abteilung IV",
        "Urteil vom 12. September 2024",
        "Gegenstand: Asylrecht, Wegweisung",
        "Regeste: Anforderungen an die Glaubhaftmachung von Verfolgung im Asylverfahren.",
        "Das Bundesverwaltungsgericht praezisiert die Anforderungen an die Glaubhaftmachung gemaess Art. 7 AsylG.",
        "Allgemeine Ausfuehrungen zu den Aspekten der Laendersituation, individuellen Gefaehrdung und Glaubwuerdigkeit der Asylvorbringen."
    ])

make_docx("Strafbefehl_Verkehrsdelikt.docx",
    "Strafbefehl",
    [
        "Staatsanwaltschaft des Kantons Zürich",
        "Strafbefehl vom 20. Februar 2026",
        "Beschuldigter: [Name anonymisiert]",
        "Tatbestand: Geschwindigkeitsüberschreitung innerorts gemäss Art. 90 Abs. 2 SVG",
        "Am 15.01.2026 um 23:42 Uhr wurde der Beschuldigte auf der Badenerstrasse, Zürich, mit einer Geschwindigkeit von 78 km/h (nach Abzug der Toleranz) gemessen. Die zulässige Höchstgeschwindigkeit beträgt 50 km/h.",
        "Strafmass: Geldstrafe von 30 Tagessätzen à CHF 150.–, bedingt, Probezeit 2 Jahre. Busse CHF 1'000.–.",
        "Einsprache: Innert 10 Tagen seit Zustellung bei der Staatsanwaltschaft."
    ])


# ============================================================
# FACHARTIKEL & PAPERS (8 Dateien)
# ============================================================

make_pdf("Paper_AI_Regulation_EU.pdf",
    "The EU AI Act: Implications for the Legal Profession",
    [
        "Authors: Prof. Dr. Sarah Chen, University of Zurich / Dr. Marco Rossi, ETH Zurich",
        "Published in: European Journal of Law and Technology, Vol. 15, No. 2 (2025)",
        "Abstract: This paper examines the impact of the EU Artificial Intelligence Act (Regulation 2024/1689) on the legal profession, with a particular focus on AI systems used for legal research, contract analysis, and predictive analytics.",
        "1. Introduction: The EU AI Act establishes a risk-based regulatory framework for AI systems. Legal technology applications may fall under various risk categories.",
        "2. Classification of Legal AI: Legal research tools are generally classified as limited risk, while AI systems used in judicial decision-making may be classified as high-risk.",
        "3. Compliance Requirements: High-risk AI systems must meet requirements regarding data governance, transparency, human oversight, and accuracy.",
        "4. Impact on Swiss Law Firms: Although Switzerland is not an EU member state, Swiss firms advising EU clients must ensure compliance.",
        "5. Conclusion: The legal profession must proactively adapt to the new regulatory landscape while leveraging the benefits of AI technology."
    ])

make_docx("Aufsatz_nDSG_Praxishinweise.docx",
    "Das neue Datenschutzgesetz (nDSG) – Praxishinweise für Unternehmen",
    [
        "Von Dr. iur. Claudia Baumgartner, Rechtsanwältin, Zürich",
        "Erschienen in: Anwaltsrevue 4/2025",
        "Das am 1. September 2023 in Kraft getretene totalrevidierte Datenschutzgesetz (nDSG) bringt erhebliche Änderungen für Unternehmen in der Schweiz.",
        "1. Informationspflicht (Art. 19 nDSG): Unternehmen müssen betroffene Personen über jede Beschaffung von Personendaten informieren. Dies gilt auch für Daten, die nicht direkt bei der betroffenen Person erhoben werden.",
        "2. Datenschutz-Folgenabschätzung (Art. 22 nDSG): Bei Datenbearbeitungen mit hohem Risiko ist vorgängig eine DSFA durchzuführen.",
        "3. Meldepflicht bei Datensicherheitsverletzungen (Art. 24 nDSG): Verletzungen der Datensicherheit sind dem EDÖB unverzüglich zu melden.",
        "4. Sanktionen (Art. 60 ff. nDSG): Natürliche Personen können mit Bussen bis CHF 250'000.– bestraft werden.",
        "Empfehlung: Unternehmen sollten ihre Datenschutzprozesse umgehend überprüfen und anpassen."
    ])

make_pdf("studie_legal_tech_schweiz_2025.pdf",
    "Legal Tech in der Schweiz - Marktueberblick 2025",
    [
        "Herausgeber: Schweizerischer Anwaltsverband (SAV) / Universitaet Luzern",
        "Executive Summary:",
        "Die vorliegende Studie untersucht den Einsatz von Legal-Tech-Loesungen in Schweizer Anwaltskanzleien und Rechtsabteilungen.",
        "Kernergebnisse:",
        "- 67% der befragten Kanzleien setzen mindestens ein Legal-Tech-Tool ein",
        "- Am haeufigsten genutzt: Dokumentenmanagement (82%), Zeiterfassung (76%), Vertragsanalyse (34%)",
        "- KI-basierte Tools: 23% nutzen KI fuer Recherche, 12% fuer Vertragsanalyse",
        "- Groesste Huerde: Datenschutzbedenken (71%), Kosten (58%), fehlende Digitalkompetenz (45%)",
        "- 89% erwarten, dass KI die juristische Arbeit in den naechsten 5 Jahren grundlegend veraendern wird"
    ])

make_pdf("article_droit_numerique_2025.pdf",
    "L'intelligence artificielle dans la pratique juridique suisse",
    [
        "Par Me Jean-Philippe Durand, Avocat, Geneve",
        "Publie dans: Revue suisse de jurisprudence (RSJ), 2025, p. 345-362",
        "Resume: Cet article analyse l'impact de l'intelligence artificielle sur la profession d'avocat en Suisse.",
        "1. Introduction: L'IA transforme rapidement le paysage juridique mondial. La Suisse n'echappe pas a cette tendance.",
        "2. Cadre reglementaire: En l'absence d'une legislation specifique sur l'IA, le droit suisse offre un cadre flexible mais incertain.",
        "3. Applications pratiques: Recherche juridique automatisee, analyse contractuelle, prediction de jugements.",
        "4. Questions ethiques: Secret professionnel, responsabilite en cas d'erreur de l'IA, biais algorithmiques.",
        "5. Conclusion: La profession juridique doit s'adapter tout en preservant les valeurs fondamentales de l'etat de droit."
    ])

make_docx("Kommentar_Art697_OR.docx",
    "Kurzkommentar zu Art. 697 OR (Auskunftsrecht)",
    [
        "Bearbeitet von Prof. Dr. Hans Keller, Universität St. Gallen",
        "Aus: Basler Kommentar, Obligationenrecht II, 6. Aufl. 2024",
        "Art. 697 OR – Auskunftsrecht",
        "I. Normzweck: Das Auskunftsrecht dient der informierten Ausübung der Aktionärsrechte, insbesondere des Stimmrechts an der Generalversammlung.",
        "II. Voraussetzungen: Der Aktionär muss darlegen, inwiefern die verlangte Auskunft für die Ausübung seiner Rechte erforderlich ist.",
        "III. Grenzen: Die Gesellschaft kann die Auskunft verweigern, soweit Geschäftsgeheimnisse oder andere schutzwürdige Interessen gefährdet würden (Art. 697 Abs. 3 OR).",
        "IV. Aktienrechtsrevision 2023: Die Revision hat das Auskunftsrecht gestärkt und den Informationsanspruch erweitert.",
        "V. Praxishinweis: Das Bundesgericht hat die Anforderungen an die Verweigerung der Auskunft in BGE 148 II 137 präzisiert."
    ])

make_pdf("whitepaper_blockchain_smart_contracts.pdf",
    "Smart Contracts and Swiss Contract Law",
    [
        "Whitepaper by the Swiss Blockchain Federation, 2025",
        "1. Executive Summary: This whitepaper examines the legal status of smart contracts under Swiss law.",
        "2. Definition: A smart contract is a self-executing program that automatically enforces the terms of an agreement.",
        "3. Legal Classification: Under Swiss law, smart contracts can constitute valid contracts if the requirements of Art. 1 ff. CO are met.",
        "4. Form Requirements: Swiss law generally follows the principle of freedom of form (Art. 11 CO). Smart contracts satisfy this requirement.",
        "5. DLT Act: The Federal Act on the Adaptation of Federal Law to Developments in Distributed Ledger Technology (DLT Act) provides a framework.",
        "6. Challenges: Immutability vs. right to rectification, jurisdiction, applicable law.",
        "7. Recommendations: Legal practitioners should advise clients on the specific risks and opportunities of smart contracts."
    ])

make_docx("Seminararbeit_KI_Haftung.docx",
    "Haftung für Künstliche Intelligenz im Schweizer Recht",
    [
        "Seminararbeit von Lena Fischer, MLaw",
        "Universität Zürich, Rechtswissenschaftliche Fakultät, HS 2025",
        "Betreuer: Prof. Dr. Michael Widmer",
        "Inhaltsverzeichnis: 1. Einleitung – 2. Haftungsgrundlagen – 3. Produktehaftung – 4. Verschuldenshaftung – 5. Kausalhaftung – 6. Rechtsvergleich EU – 7. Fazit",
        "Abstract: Diese Arbeit untersucht, wie das geltende Schweizer Haftungsrecht auf Schäden durch KI-Systeme angewendet werden kann.",
        "Zentrale These: Das bestehende Haftungsrecht ist grundsätzlich geeignet, KI-bezogene Schäden zu erfassen, jedoch bestehen Lücken bei autonomen Systemen.",
        "Fazit: Es wird empfohlen, eine spezifische Kausalhaftung für den Betrieb von Hochrisiko-KI-Systemen einzuführen, analog zur Haftung des Tierhalters (Art. 56 OR)."
    ])

make_pdf("newsletter_compliance_q1_2026.pdf",
    "Compliance Newsletter Q1/2026",
    [
        "Herausgeber: Compliance Solutions AG, Zuerich",
        "Themenueberblick:",
        "1. Neuerungen im Geldwaeschereigesetz (GwG): Die FINMA hat neue Rundschreiben zur Sorgfaltspflicht veroeffentlicht.",
        "2. ESG-Reporting: Ab 2026 gelten erweiterte Berichtspflichten fuer grosse Unternehmen (Art. 964a ff. OR).",
        "3. Whistleblowing: Der Bundesrat hat einen Vorentwurf fuer ein Whistleblower-Schutzgesetz in die Vernehmlassung geschickt.",
        "4. Sanktionen: Aktuelle Entwicklungen bei den Sanktionen gegen Russland.",
        "5. Veranstaltungshinweis: Compliance Summit 2026, 15. Mai, Kongresshaus Zuerich."
    ])


# ============================================================
# KORRESPONDENZ & NOTIZEN (10 Dateien)
# ============================================================

make_txt("Meeting_Notes_15.03.txt",
    """Meeting Notes – Projekt Alpha Statusmeeting
Datum: 15.03.2026, 14:00-15:30
Teilnehmer: Dr. Keller, Frau Meier, Herr Brunner, Frau Rossi

Agenda:
1. Stand Vertragsverhandlungen
   - NDA unterzeichnet
   - SPA Entwurf v2 liegt vor, Kommentare bis Freitag
   - Due Diligence Dataroom wird nächste Woche eröffnet

2. Offene Punkte
   - Steuergutachten noch ausstehend (Deadline: 22.03.)
   - Bewertung der IP-Rechte unklar → externen Gutachter beiziehen
   - Arbeitnehmervertretung muss informiert werden

3. Next Steps
   - Frau Meier: Dataroom-Struktur bis Montag
   - Herr Brunner: Steuergutachten nachfassen
   - Dr. Keller: Gespräch mit Gewerkschaft terminieren

Nächstes Meeting: 22.03.2026, 10:00""")

make_docx("Brief_an_Gegenpartei_Mahnung.docx",
    "Mahnung – Zahlungsaufforderung",
    [
        "Kanzlei Weber & Partner\nTalstrasse 18\n8001 Zürich\n\nEinschreiben",
        "An: InnoSoft GmbH\nMaximilianstrasse 42\n80538 München\nDeutschland",
        "Zürich, 10. März 2026",
        "Betreff: Mahnung – Ausstehende Zahlung gemäss Urteil HG250124",
        "Sehr geehrte Damen und Herren,",
        "namens und im Auftrag unserer Mandantin, der TechVision AG, Zürich, fordern wir Sie hiermit letztmalig auf, den rechtskräftig zugesprochenen Betrag von CHF 150'000.– innert 10 Tagen auf das nachfolgende Konto zu überweisen.",
        "IBAN: CH93 0076 2011 6238 5295 7\nBank: Zürcher Kantonalbank\nKontoinhaber: TechVision AG",
        "Sollte die Zahlung nicht fristgerecht eingehen, werden wir ohne weitere Ankündigung die Zwangsvollstreckung einleiten.",
        "Freundliche Grüsse\nKanzlei Weber & Partner\n\nDr. iur. Stefan Weber\nRechtsanwalt"
    ])

make_txt("notiz_telefonat_mandant.txt",
    """Telefonnotiz
Datum: 12.03.2026, 16:45
Mandant: Hr. Baumgartner (DataFlow Systems AG)
Gesprächspartner: RA Dr. Weber

Inhalt:
- Mandant besorgt wegen EDÖB-Empfehlung zu KI-Tools
- Nutzt aktuell ChatGPT und Claude für interne Recherche
- Fragt: Müssen wir DSFA machen?
- Fragt: Können Mitarbeitende persönlich gebüsst werden?
- Will Gutachten zum Thema KI und Berufsgeheimnis

→ Bis nächste Woche Kurzgutachten erstellen (Budget: 5h)
→ Termin für Besprechung: 20.03., 10:00""")

make_docx("memo_intern_due_diligence.docx",
    "Internes Memo – Due Diligence MedTech Innovations",
    [
        "VERTRAULICH – ANWALTSPRIVILEG",
        "Von: Ass.-jur. Laura Schmidt\nAn: Dr. Stefan Weber\nDatum: 08.03.2026",
        "Betreff: Erste Erkenntnisse Due Diligence MedTech Innovations AG",
        "Zusammenfassung der bisherigen Prüfung:",
        "1. Gesellschaftsrechtlich: Kapitalstruktur sauber, keine Vinkulierungsklauseln die problematisch wären.",
        "2. Arbeitsrechtlich: 3 Schlüsselmitarbeiter ohne Konkurrenzverbot → Retention-Vereinbarungen empfohlen.",
        "3. IP: 12 Patente, davon 2 in Lizenzstreit (DE) → Risiko ca. EUR 2 Mio.",
        "4. Regulatorisch: CE-Zertifizierungen aktuell, MDR-Konformität bestätigt.",
        "5. Datenschutz: Keine DSFA vorhanden trotz Verarbeitung von Patientendaten → Red Flag!",
        "Empfehlung: Kaufpreisreduktion von CHF 2-3 Mio. wegen IP-Risiken. DSFA-Nachholung als Closing Condition."
    ])

make_txt("todo_liste_KW11.txt",
    """To-Do Liste KW 11/2026 – RA Dr. Weber

[ ] SPA MedTech: Kommentare von Partner Brunner einarbeiten
[ ] Mahnung InnoSoft: Einschreiben versenden
[x] NDA Projekt Alpha: Unterzeichnete Fassung archivieren
[ ] Kurzgutachten KI/Berufsgeheimnis für Hr. Baumgartner
[ ] EDÖB-Empfehlung: Stellungnahme vorbereiten (Frist: 15.04.)
[ ] Seminar Datenschutz: Slides aktualisieren bis 20.03.
[ ] Mandatsabrechnung Q1: Bis Ende März
[ ] Neue Praktikantin: Einführung am Montag
[x] Mietvertrag Talstrasse: Verlängerung unterschrieben
[ ] Compliance Newsletter lesen""")

make_docx("email_print_mandantenanfrage.docx",
    "E-Mail-Ausdruck: Mandantenanfrage",
    [
        "Von: h.baumgartner@dataflow.ch\nAn: weber@kanzlei-weber.ch\nDatum: 05.03.2026, 09:12\nBetreff: Dringend – KI-Policy für unser Unternehmen",
        "Lieber Herr Dr. Weber,",
        "wie heute Morgen telefonisch besprochen, benötigen wir dringend Unterstützung bei der Erstellung einer unternehmensweiten KI-Nutzungsrichtlinie.",
        "Hintergrund: Wir haben festgestellt, dass verschiedene Abteilungen unterschiedliche KI-Tools nutzen (ChatGPT, Claude, Copilot). Es gibt keine einheitliche Regelung bezüglich Datenschutz, Vertraulichkeit und Qualitätssicherung.",
        "Können Sie uns ein Angebot machen für:\n1. Analyse des Ist-Zustands\n2. Erstellung einer KI-Nutzungsrichtlinie\n3. Schulung der Abteilungsleiter (ca. 15 Personen)",
        "Budget: Wir haben intern CHF 30'000 – 40'000 genehmigt.",
        "Zeitrahmen: Möglichst bis Ende April.",
        "Freundliche Grüsse\nHans Baumgartner\nCEO, DataFlow Systems AG"
    ])

make_txt("ideen_blogpost_legaltech.txt",
    """Ideen für Blogpost: "KI in der Anwaltskanzlei – Was geht, was nicht?"

Zielgruppe: Anwälte und Kanzleimitarbeitende

Gliederung:
- Einleitung: Status quo in CH Kanzleien (Zahlen aus SAV-Studie)
- Was KI heute schon kann: Recherche, Vertragsanalyse, Übersetzung
- Was KI (noch) nicht kann: Strategische Beratung, Verhandlung, Empathie
- Rechtliche Grenzen: Berufsgeheimnis, Datenschutz, Sorgfaltspflicht
- Praktische Tipps: Do's and Don'ts
- Fazit: KI als Werkzeug, nicht als Ersatz

Quellen:
- SAV Legal Tech Studie 2025
- EDÖB Empfehlung
- EU AI Act
- BGE 148 II 137 (Auskunftsrecht analog?)

Deadline: Ende März für Website
Wörter: ca. 1500""")

make_docx("gutachten_entwurf_ki_berufsgeheimnis.docx",
    "ENTWURF – Kurzgutachten: Einsatz von KI-Tools und anwaltliches Berufsgeheimnis",
    [
        "Erstellt für: DataFlow Systems AG / z.Hd. CEO Hans Baumgartner",
        "Erstellt von: Kanzlei Weber & Partner, Dr. Stefan Weber",
        "Datum: ENTWURF – 13.03.2026",
        "I. Fragestellung: Dürfen Mitarbeitende der DataFlow Systems AG KI-Tools wie ChatGPT oder Claude nutzen, ohne das Berufsgeheimnis zu verletzen?",
        "II. Rechtslage:",
        "Art. 321 StGB schützt das Berufsgeheimnis. Die Weitergabe von Mandantendaten an KI-Anbieter könnte eine strafbare Offenbarung darstellen.",
        "Das nDSG verlangt eine Rechtsgrundlage für jede Bearbeitung von Personendaten (Art. 6 nDSG).",
        "III. Beurteilung:",
        "Die Nutzung von KI-Tools ist grundsätzlich zulässig, sofern keine Personendaten oder vertraulichen Mandantsinformationen eingegeben werden.",
        "IV. Empfehlungen: [WIRD NOCH ERGÄNZT]",
        "V. Vorbehalt: Dieses Gutachten ist ein Entwurf und noch nicht freigegeben."
    ])

make_txt("spesen_maerz_2026.txt",
    """Spesenabrechnung März 2026 – Dr. Stefan Weber

05.03. Mandantenessen mit Hr. Baumgartner (Kronenhalle) CHF 245.50
07.03. Taxi zum Handelsgericht CHF 32.00
10.03. Zugfahrt Zürich-Bern (SPA-Verhandlung) CHF 51.00
10.03. Mittagessen Bern CHF 38.50
12.03. Parkgebühren Flughafen (Mandantenbesuch) CHF 45.00
15.03. Fachliteratur: BSK OR II, Neuauflage CHF 298.00
18.03. Konferenzteilnahme Legal Tech Summit CHF 450.00

Total: CHF 1'160.00

Belege: Beigelegt / gescannt""")

make_docx("protokoll_generalversammlung.docx",
    "Protokoll der ordentlichen Generalversammlung",
    [
        "TechVision AG, Zürich",
        "Datum: 28. Februar 2026, 14:00 Uhr\nOrt: Hotel Baur au Lac, Zürich",
        "Anwesend: 85% des Aktienkapitals vertreten",
        "Vorsitz: Dr. Maria Schneider, Verwaltungsratspräsidentin",
        "Traktanden:",
        "1. Genehmigung des Geschäftsberichts 2025: Einstimmig genehmigt.",
        "2. Genehmigung der Jahresrechnung 2025: Umsatz CHF 42 Mio. (+12%), EBIT CHF 6.3 Mio. Einstimmig genehmigt.",
        "3. Verwendung des Bilanzgewinns: Dividende CHF 2.50 pro Aktie. Mit 92% der Stimmen genehmigt.",
        "4. Entlastung des Verwaltungsrates: Einstimmig erteilt.",
        "5. Wahl Revisionsstelle: PwC Zürich für ein weiteres Jahr gewählt.",
        "6. Verschiedenes: Keine Wortmeldungen.",
        "Protokollführer: Dr. Stefan Weber, Kanzlei Weber & Partner"
    ])


# ============================================================
# RECHNUNGEN & FINANZEN (6 Dateien)
# ============================================================

make_pdf("Rechnung_2024-0847.pdf",
    "Rechnung Nr. 2024-0847",
    [
        "Kanzlei Weber & Partner\nTalstrasse 18, 8001 Zuerich\nMWST-Nr.: CHE-123.456.789",
        "An: TechVision AG, Zuerich",
        "Datum: 28.02.2026",
        "Betreff: Honorar fuer Rechtsberatung Januar-Februar 2026",
        "Leistungsbeschreibung:",
        "- Beratung Vertragsrecht (SPA MedTech): 42.5 Std. x CHF 450 = CHF 19'125.00",
        "- Due Diligence Begleitung: 28.0 Std. x CHF 450 = CHF 12'600.00",
        "- Korrespondenz und Verhandlungen: 15.5 Std. x CHF 350 = CHF 5'425.00",
        "Zwischensumme: CHF 37'150.00",
        "Auslagen (Handelsregister, Kopien): CHF 385.00",
        "Total exkl. MWST: CHF 37'535.00",
        "MWST 8.1%: CHF 3'040.34",
        "Total inkl. MWST: CHF 40'575.34",
        "Zahlbar innert 30 Tagen auf IBAN CH93 0076 2011 6238 5295 7"
    ])

make_pdf("offerte_ki_policy_dataflow.pdf",
    "Offerte: KI-Nutzungsrichtlinie fuer DataFlow Systems AG",
    [
        "Kanzlei Weber & Partner, Zuerich",
        "Offerte vom 13.03.2026",
        "Adressat: DataFlow Systems AG, z.Hd. Hr. Hans Baumgartner",
        "1. Analyse Ist-Zustand: 12 Std. x CHF 350 = CHF 4'200",
        "2. Erstellung KI-Nutzungsrichtlinie: 25 Std. x CHF 350 = CHF 8'750",
        "3. Datenschutz-Folgenabschaetzung: 15 Std. x CHF 350 = CHF 5'250",
        "4. Schulung Abteilungsleiter (1 Tag): Pauschal CHF 4'500",
        "5. Projektleitung und Koordination: 8 Std. x CHF 450 = CHF 3'600",
        "Total exkl. MWST: CHF 26'300",
        "MWST 8.1%: CHF 2'130.30",
        "Total inkl. MWST: CHF 28'430.30",
        "Gueltig bis: 15.04.2026"
    ])

make_pdf("honorarnote_2026_Q1_Mueller.pdf",
    "Honorarnote Q1/2026",
    [
        "Kanzlei Weber & Partner, Zuerich",
        "Mandant: Thomas Mueller (Aufhebungsvereinbarung)",
        "Zeitraum: Januar - Maerz 2026",
        "Leistungen:",
        "- Erstberatung und Sachverhaltsanalyse: 3.0 Std.",
        "- Verhandlung mit Arbeitgeberin: 5.5 Std.",
        "- Entwurf Aufhebungsvereinbarung: 4.0 Std.",
        "- Korrespondenz: 2.5 Std.",
        "Total: 15.0 Std. x CHF 400 = CHF 6'000.00",
        "MWST 8.1%: CHF 486.00",
        "Total: CHF 6'486.00",
        "Zahlbar innert 30 Tagen."
    ])

make_pdf("budget_2026_kanzlei.pdf",
    "Budget 2026 - Kanzlei Weber & Partner",
    [
        "VERTRAULICH",
        "Erstellt: Januar 2026",
        "Einnahmen:",
        "- Honorare Mandate: CHF 1'850'000 (Vorjahr: CHF 1'620'000)",
        "- Pauschalmandate: CHF 420'000",
        "- Diverse: CHF 30'000",
        "Total Einnahmen: CHF 2'300'000",
        "",
        "Ausgaben:",
        "- Personalkosten (6 Anwaelte, 3 Support): CHF 1'380'000",
        "- Miete Talstrasse: CHF 115'200",
        "- IT und Software: CHF 85'000 (inkl. Legal Tech Tools)",
        "- Versicherungen: CHF 42'000",
        "- Marketing: CHF 35'000",
        "- Weiterbildung: CHF 28'000",
        "- Uebrige: CHF 64'800",
        "Total Ausgaben: CHF 1'750'000",
        "",
        "Budgetierter Gewinn: CHF 550'000"
    ])

make_txt("zeiterfassung_kw10.txt",
    """Zeiterfassung KW 10/2026 – Dr. Stefan Weber

Montag 03.03.
  TechVision/MedTech DD    4.5h    Due Diligence Review
  DataFlow/Allgemein       1.0h    Telefonat Baumgartner
  Intern                   0.5h    Teamsitzung

Dienstag 04.03.
  TechVision/MedTech DD    6.0h    Dataroom Review
  Müller/Arbeitsrecht      1.5h    Verhandlung Aufhebung

Mittwoch 05.03.
  TechVision/MedTech DD    3.0h    SPA Kommentare
  DataFlow/KI-Policy       2.0h    Recherche nDSG/KI
  TechVision/InnoSoft      2.0h    Mahnung vorbereiten

Donnerstag 06.03.
  TechVision/MedTech DD    5.5h    Vertragsverhandlung Bern
  Reisezeit                2.0h    Zürich-Bern retour

Freitag 07.03.
  TechVision/MedTech DD    3.0h    Memo Due Diligence
  Intern                   2.0h    Weiterbildung Legal Tech
  Admin                    1.0h    Rechnungen, Spesen

Total Woche: 34.0h (fakturierbar: 28.5h)""")

make_pdf("Rechnung_CloudSecure_hosting.pdf",
    "Invoice / Rechnung",
    [
        "CloudSecure AG, Zurich",
        "Invoice No.: CS-2026-0312",
        "Date: 01.03.2026",
        "To: Kanzlei Weber & Partner, Zurich",
        "Service Period: March 2026",
        "Description:",
        "- Managed Cloud Hosting (Legal DMS): CHF 1'200.00",
        "- Additional Storage (500 GB): CHF 150.00",
        "- SSL Certificate Renewal: CHF 85.00",
        "- 24/7 Support Package: CHF 350.00",
        "Subtotal: CHF 1'785.00",
        "VAT 8.1%: CHF 144.59",
        "Total: CHF 1'929.59",
        "Payment Terms: 30 days net",
        "IBAN: CH12 3456 7890 1234 5678 9"
    ])


# ============================================================
# PRÄSENTATIONEN (4 Dateien)
# ============================================================

make_pptx("Präsentation_Mandantenmeeting.pptx",
    "Projekt Alpha – Statusbericht",
    [
        ("Projektübersicht", "• Akquisition MedTech Innovations AG\n• Kaufpreis: CHF 45 Mio.\n• Closing geplant: Q2 2026\n• Status: Due Diligence Phase"),
        ("Due Diligence – Zusammenfassung", "• Gesellschaftsrecht: ✓ OK\n• Arbeitsrecht: ⚠ Retention Agreements nötig\n• IP: ⚠ 2 Patente in Streit\n• Datenschutz: ✗ Keine DSFA vorhanden\n• Regulatorisch: ✓ CE/MDR konform"),
        ("Nächste Schritte", "1. SPA Entwurf finalisieren (bis 22.03.)\n2. Steuergutachten einholen\n3. Retention Agreements verhandeln\n4. DSFA als Closing Condition\n5. Signing geplant: April 2026"),
    ])

make_pptx("slides_datenschutz_schulung.pptx",
    "Datenschutz-Schulung für Mitarbeitende",
    [
        ("Warum Datenschutz?", "• Neues Datenschutzgesetz (nDSG) seit 01.09.2023\n• Bussen bis CHF 250'000 für natürliche Personen\n• Reputationsrisiko bei Datenpannen\n• Vertrauen der Kunden und Mandanten"),
        ("Was sind Personendaten?", "• Name, Adresse, E-Mail\n• Telefonnummer\n• AHV-Nummer\n• Gesundheitsdaten (besonders schützenswert!)\n• Finanzdaten"),
        ("Die 7 Grundsätze", "1. Rechtmässigkeit\n2. Treu und Glauben\n3. Verhältnismässigkeit\n4. Zweckbindung\n5. Richtigkeit\n6. Datensicherheit\n7. Transparenz"),
        ("Do's and Don'ts", "DO: Bildschirm sperren, Dokumente shredden, verschlüsselt mailen\nDON'T: Passwörter teilen, Daten per WhatsApp senden, USB-Sticks unverschlüsselt nutzen"),
    ])

make_pptx("pitch_deck_legaltech_tool.pptx",
    "JurisAssist Pro – KI-Rechtsrecherche",
    [
        ("Das Problem", "• Juristische Recherche dauert Stunden\n• Relevante Urteile werden übersehen\n• Hohe Kosten für Mandanten\n• Junior Associates verbringen 60% ihrer Zeit mit Recherche"),
        ("Unsere Lösung", "• KI-gestützte Recherche in Sekunden\n• Über 500'000 Schweizer Urteile indexiert\n• Automatische Relevanz-Bewertung\n• Integriert in bestehende DMS-Systeme"),
        ("Markt & Traction", "• 12'000+ Anwälte in der Schweiz\n• TAM: CHF 120 Mio. p.a.\n• 45 Kanzleien in Pilotphase\n• NPS Score: 72"),
        ("Ask", "• Series A: CHF 8 Mio.\n• Verwendung: Produktentwicklung (50%), Sales (30%), Team (20%)\n• Break-even: Q4 2027"),
    ])

make_pptx("vortrag_aktienrechtsrevision.pptx",
    "Die Aktienrechtsrevision 2023 – Ein Jahr danach",
    [
        ("Überblick der Änderungen", "• Kapitalband (Art. 653s ff. OR)\n• Virtuelle Generalversammlung (Art. 701d OR)\n• Geschlechterrichtwerte (Art. 734f OR)\n• Erweitertes Auskunftsrecht (Art. 697 OR)\n• Transparenz Rohstoffsektor"),
        ("Praxiserfahrungen", "• Kapitalband: Noch wenig genutzt (ca. 120 Gesellschaften)\n• Virtuelle GV: Grosser Zuspruch, rechtliche Fragen offen\n• Geschlechterrichtwerte: 65% der SMI-Gesellschaften erfüllt\n• Auskunftsrecht: BGE 148 II 137 als Leitentscheid"),
        ("Offene Fragen", "• Hybride GV: Stimmrechtsausübung bei technischen Problemen?\n• Kapitalband: Steuerliche Behandlung?\n• Auskunftsrecht: Wo liegt die Grenze zum Missbrauch?\n• Digitale Aktien: Praxistauglichkeit?"),
    ])


# ============================================================
# SONSTIGE DATEIEN (4 Dateien)
# ============================================================

make_txt("passwörter_NICHT_LÖSCHEN.txt",
    """ACHTUNG: Diese Datei sollte NICHT in einem unverschlüsselten Ordner liegen!

Zugangsdaten (DEMO-DATEN – nicht echt):
- DMS System: admin / Kanzlei2026!
- CloudSecure Portal: weber@kanzlei-weber.ch / ********
- Handelsregister Online: User4711 / ********
- EDÖB Portal: kanzlei.weber / ********

BITTE IN PASSWORT-MANAGER ÜBERTRAGEN UND DIESE DATEI LÖSCHEN!""")

make_txt("scan0042.txt",
    """[OCR-Scan – Qualität: mittel]

Einschreiben
Betreibungsamt Zürich 1
Datum: 02.03.2026

Zahlungsbefehl Nr. 2026-ZH1-45892

Gläubigerin: TechVision AG, Zürich
Schuldnerin: InnoSoft GmbH, München

Forderung: CHF 150'000.00
Zins: 5% seit 15.01.2026
Kosten: CHF 103.00

Rechtsvorschlag: Frist 10 Tage ab Zustellung""")

make_docx("lebenslauf_praktikantin_neu.docx",
    "Lebenslauf – Anna-Lisa Berger",
    [
        "Anna-Lisa Berger\nSeefeldstrasse 123, 8008 Zürich\nanna.berger@uzh.ch\nTel: +41 78 123 45 67",
        "Ausbildung:\n- MLaw, Universität Zürich (2024), Note: 5.4\n- BLaw, Universität Zürich (2022), Note: 5.2\n- Matura, Kantonsschule Hottingen (2019)",
        "Praktika:\n- 6 Monate Bezirksgericht Zürich (2024)\n- 3 Monate Kanzlei Lenz & Staehelin, Zürich (2023)",
        "Sprachkenntnisse:\n- Deutsch (Muttersprache)\n- Französisch (C1)\n- Englisch (C2)\n- Italienisch (B1)",
        "Schwerpunkte: Gesellschaftsrecht, M&A, Datenschutzrecht",
        "Referenzen: Auf Anfrage"
    ])

make_txt("notiz_parkplatz_kündigung.txt",
    """Notiz: Parkplatz Talstrasse

Der Vermieter hat mitgeteilt, dass die Parkplätze in der
Tiefgarage ab Juli 2026 neu CHF 350/Monat kosten (bisher 280).

Optionen:
1. Neuen Preis akzeptieren (3 PP x 350 = 1050/Mt)
2. Auf 2 Parkplätze reduzieren
3. Alternative Garage suchen (Parkhaus Hohe Promenade?)

→ An Teamsitzung besprechen""")


# ============================================================
# FERTIG
# ============================================================

files = os.listdir(OUT)
print(f"\n✓ {len(files)} Dateien erstellt in {OUT}/\n")
for f in sorted(files):
    size = os.path.getsize(os.path.join(OUT, f))
    print(f"  {f:55s} {size:>8,} Bytes")
